"""Create editable PPTX and DOCX deliverables for SeatSense AI."""

from __future__ import annotations

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from src.predict import ensure_project_ready  # noqa: E402


DELIVERABLES_DIR = PROJECT_ROOT / "deliverables"
OUTPUTS_DIR = PROJECT_ROOT / "outputs"


SLIDES = [
    (
        "SeatSense AI",
        [
            "Dynamic Ticket Pricing & Demand Intelligence for Sports Venues",
            "Enterprise revenue intelligence for ticketing teams and venue operators",
            "Predict demand, recommend guarded prices, and explain approval risk.",
        ],
    ),
    (
        "Business Problem",
        [
            "Static prices miss fast-changing demand signals.",
            "Low-demand games leave seats unsold.",
            "High-demand games are often underpriced.",
            "Secondary markets capture value that should stay with the venue.",
        ],
    ),
    (
        "Market Context and Existing Customers",
        [
            "Customers: teams, arenas, venues, and ticketing departments.",
            "Users: revenue managers, operations, marketing, fan experience, executives.",
            "Existing market already values dynamic pricing tools.",
            "SeatSense adds explainability, guardrails, and business KPI reporting.",
        ],
    ),
    (
        "Who, Why, How Framework",
        [
            "Who: teams, venues, fans, executives, and ticketing staff.",
            "Why: demand is multivariate, volatile, and financially material.",
            "How: model prediction plus governed pricing engine plus explanation.",
            "Human approval remains in the decision workflow.",
        ],
    ),
    (
        "Why AI Is the Right Solution",
        [
            "Rules miss interactions among opponent, weather, sentiment, and resale gap.",
            "Classification predicts Low, Medium, or High demand.",
            "Probabilities support confidence-aware pricing decisions.",
            "OpenAI-enabled explanation translates outputs into revenue-manager language with safe fallback.",
        ],
    ),
    (
        "Product Overview: SeatSense AI",
        [
            "Use sample data or map client CSV exports.",
            "Predict demand tier and class probability.",
            "Recommend guarded price changes.",
            "Estimate revenue uplift and resale leakage reduction.",
            "Generate executive-ready explanations from model output.",
            "Flag decisions that require human approval.",
        ],
    ),
    (
        "Data Pipeline and Feature Design",
        [
            "Semi-synthetic 20,500-row game-section-pricing-window dataset.",
            "Five seasons, 410 games, 10 seat sections, five pricing windows.",
            "Features include current inventory, sell-through, traffic, search, social, weather, price, resale, section, and historical behavior.",
            "Leakage audit excludes final attendance, realized revenue, target labels, oracle price, and fairness-only affordability columns.",
            "Ticketing-agency adapter maps vendor feeds, then Pydantic validates the SeatSense schema before retraining.",
        ],
    ),
    (
        "Model Architecture",
        [
            "Demand classifier: simple decision-tree baseline vs tuned random forest.",
            "Sell-through regressor estimates final inventory conversion.",
            "Scalper-risk classifier estimates resale leakage pressure.",
            "Time-based split: first 3 seasons train, 4th validation, 5th test.",
            "Pricing optimizer simulates candidate prices and applies guardrails.",
        ],
    ),
    (
        "Model Performance",
        [
            "Final demand accuracy is about 86.7% with macro F1 about 86.1%.",
            "High-demand recall is about 89.8%; ROC-AUC OVR is about 96.9%.",
            "Sell-through RMSE is about 0.059; scalper-risk macro F1 is about 88.2%.",
            "Confusion matrix shows costly pricing mistakes.",
            "Leakage audit confirms no post-event or affordability-only columns are used as model inputs.",
            "Data quality widgets include demand mix, price gap scatter, signal bubble, and section drill-down.",
        ],
    ),
    (
        "Pricing Recommendation Workflow",
        [
            "Inputs: current price, predicted demand, resale gap, days before game, section, sell-through.",
            "Optional live APIs enrich ticket marketplace, weather, sports, social, and traffic signals.",
            "Apply price caps, affordability rules, and premium-game approval.",
            "Estimate revenue uplift and leakage reduction.",
            "Return business action and approval status.",
        ],
    ),
    (
        "Streamlit Product Screens",
        [
            "Executive Overview frames the problem and KPIs.",
            "Pricing Workbench shows live API signal mode, probabilities, recommended price, and approval logic.",
            "Model Performance shows flexible CSV import, agency-feed detection, schema validation, and drill-down charts.",
            "Demand Intelligence explains drivers behind the forecast.",
        ],
    ),
    (
        "Business Impact and KPIs",
        [
            "Revenue per game and total ticket revenue uplift.",
            "Sell-through and attendance rate.",
            "Secondary-market leakage reduction.",
            "Average price gap versus resale market.",
            "Fan affordability score and approval rate.",
        ],
    ),
    (
        "Responsible AI and Fairness Guardrails",
        [
            "Lower-income fans may be priced out without controls.",
            "No protected attributes are used directly.",
            "Affordable ticket reserve and value-section caps.",
            "Human review for premium games and large price moves.",
            "Bias audit by affordability segment.",
        ],
    ),
    (
        "Implementation Roadmap",
        [
            "Phase 1: historical data integration and offline validation.",
            "Phase 2: pilot with revenue-manager review.",
            "Phase 3: controlled A/B testing by section and game type.",
            "Phase 4: production monitoring for drift, fairness, and fan sentiment.",
        ],
    ),
    (
        "Conclusion and Ask",
        [
            "SeatSense captures revenue that static pricing misses.",
            "It reduces resale leakage while protecting attendance.",
            "It combines trained predictive AI with explainable recommendations and fallback behavior.",
            "Ask: approve a pilot using real historical venue data.",
        ],
    ),
]


def create_pptx() -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    navy = RGBColor(17, 24, 39)
    blue = RGBColor(37, 99, 235)
    muted = RGBColor(100, 116, 139)
    light = RGBColor(248, 250, 252)

    for idx, (title, bullets) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = navy if idx == 1 else RGBColor(255, 255, 255)

        if idx == 1:
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.8), Inches(1.2))
            title_frame = title_box.text_frame
            title_frame.clear()
            p = title_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            body = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(9.2), Inches(2.2))
            frame = body.text_frame
            frame.clear()
            for i, bullet in enumerate(bullets):
                para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                para.text = bullet
                para.font.size = Pt(22 if i == 0 else 18)
                para.font.color.rgb = RGBColor(226, 232, 240)

            accent = slide.shapes.add_shape(1, Inches(9.8), Inches(1.35), Inches(2.7), Inches(4.7))
            accent.fill.solid()
            accent.fill.fore_color.rgb = blue
            accent.line.color.rgb = blue
            tx = accent.text_frame
            tx.clear()
            p = tx.paragraphs[0]
            p.text = "AI pricing\nwith guardrails"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(24)
            p.font.bold = True
            continue

        title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.42), Inches(11.9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = navy

        rule = slide.shapes.add_shape(1, Inches(0.65), Inches(1.25), Inches(1.35), Inches(0.06))
        rule.fill.solid()
        rule.fill.fore_color.rgb = blue
        rule.line.color.rgb = blue

        body = slide.shapes.add_textbox(Inches(0.78), Inches(1.58), Inches(6.2), Inches(4.9))
        frame = body.text_frame
        frame.clear()
        frame.word_wrap = True
        for i, bullet in enumerate(bullets):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(20)
            para.font.color.rgb = navy
            para.space_after = Pt(11)

        stage = slide.shapes.add_shape(1, Inches(7.45), Inches(1.55), Inches(4.9), Inches(4.65))
        stage.fill.solid()
        stage.fill.fore_color.rgb = light
        stage.line.color.rgb = RGBColor(219, 227, 239)

        visual = slide.shapes.add_textbox(Inches(7.82), Inches(2.05), Inches(4.15), Inches(3.55))
        visual_frame = visual.text_frame
        visual_frame.clear()
        vp = visual_frame.paragraphs[0]
        vp.text = visual_label_for_slide(idx)
        vp.font.size = Pt(22)
        vp.font.bold = True
        vp.font.color.rgb = blue
        vp.alignment = PP_ALIGN.CENTER

        footer = slide.shapes.add_textbox(Inches(0.68), Inches(6.95), Inches(11.9), Inches(0.25))
        fp = footer.text_frame.paragraphs[0]
        fp.text = "SeatSense AI · Revenue Intelligence"
        fp.font.size = Pt(9)
        fp.font.color.rgb = muted

        if idx == 9:
            image_path = OUTPUTS_DIR / "confusion_matrix.png"
            if image_path.exists():
                slide.shapes.add_picture(str(image_path), Inches(7.72), Inches(1.86), width=Inches(4.0))
        elif idx == 12:
            image_path = OUTPUTS_DIR / "feature_importance.png"
            if image_path.exists():
                slide.shapes.add_picture(str(image_path), Inches(7.72), Inches(1.86), width=Inches(4.0))

    path = DELIVERABLES_DIR / "SeatSense_AI_Final_Presentation.pptx"
    prs.save(path)
    return path


def visual_label_for_slide(index: int) -> str:
    labels = {
        2: "Revenue loss\nfrom static pricing",
        3: "Stakeholder map",
        4: "Who · Why · How",
        5: "Signals to decision",
        6: "Product workflow",
        7: "Data pipeline",
        8: "Model stack",
        9: "Confusion matrix",
        10: "Pricing guardrails",
        11: "Streamlit demo",
        12: "KPI dashboard",
        13: "Responsible AI",
        14: "Roadmap",
        15: "Pilot ask",
    }
    return labels.get(index, "SeatSense AI")


def create_docx() -> Path:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    report_md = DELIVERABLES_DIR / "Project_Report.md"
    report_text = report_md.read_text(encoding="utf-8")
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.72)
    section.bottom_margin = Inches(0.72)
    section.left_margin = Inches(0.78)
    section.right_margin = Inches(0.78)

    styles = doc.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(11)
    styles["Heading 1"].font.name = "Times New Roman"
    styles["Heading 1"].font.size = Pt(16)
    styles["Heading 1"].font.color.rgb = RGBColor(17, 24, 39)
    styles["Heading 2"].font.name = "Times New Roman"
    styles["Heading 2"].font.size = Pt(13)
    styles["Heading 2"].font.color.rgb = RGBColor(29, 78, 216)

    for raw_line in report_text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("# "):
            para = doc.add_paragraph()
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = para.add_run(line[2:])
            run.bold = True
            run.font.name = "Times New Roman"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(17, 24, 39)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)

    path = DELIVERABLES_DIR / "Project_Report.docx"
    doc.save(path)
    return path


def main() -> None:
    ensure_project_ready()
    DELIVERABLES_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = create_pptx()
    docx_path = create_docx()
    print(f"PPTX created: {pptx_path}")
    print(f"DOCX created: {docx_path}")


if __name__ == "__main__":
    main()
